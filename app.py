from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.responses import FileResponse, HTMLResponse
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageOps
import io
import os
import shutil
import gc  # Importation du nettoyeur de mémoire (Garbage Collector)

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

def appliquer_filtre_image(image_blob):
    """Applique un filtre bleu Decathlon avec optimisation drastique de la RAM"""
    try:
        img = Image.open(io.BytesIO(image_blob))
        
        # 1. OPTIMISATION : Limiter la taille à du Full HD max
        img.thumbnail((1920, 1080), Image.Resampling.LANCZOS)
        
        # 2. Sécurité : Retirer la transparence avant colorisation (sinon crash)
        if img.mode in ('RGBA', 'P', 'LA'):
            img = img.convert('RGB')
            
        img = img.convert("L")
        
        # Appliquer le duotone
        img_colorized = ImageOps.colorize(
            img, 
            black=(20, 25, 60),      # Ombres
            white=(255, 255, 255),   # Lumières
            mid=(59, 73, 184)        # Tons moyens
        )
        
        out_io = io.BytesIO()
        # 3. OPTIMISATION : Sauvegarder en JPEG léger plutôt qu'en PNG
        img_colorized.save(out_io, format="JPEG", quality=85)
        
        # Libération manuelle de la mémoire
        img.close()
        img_colorized.close()
        
        return out_io
    except Exception as e:
        print(f"Image ignorée (format non supporté) : {e}")
        return None  # On renvoie None pour que le script sache qu'il faut ignorer cette image

@app.post("/api/appliquer-charte")
async def appliquer_charte(fichier: UploadFile = File(...)):
    input_path = f"temp_{fichier.filename}"
    
    try:
        with open(input_path, "wb") as buffer:
            shutil.copyfileobj(fichier.file, buffer)

        prs = Presentation(input_path)
        bleu_fond = RGBColor(59, 73, 184)
        texte_blanc = RGBColor(255, 255, 255)
        
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        for slide in prs.slides:
            # Fond bleu
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = bleu_fond

            shapes_to_delete = []

            for shape in slide.shapes:
                # --- GESTION DES IMAGES ---
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_blob = shape.image.blob
                    
                    # Tentative d'application du filtre
                    nouvelle_image_io = appliquer_filtre_image(image_blob)
                    
                    if nouvelle_image_io:  # Si l'image a été modifiée avec succès
                        left, top = shape.left, shape.top
                        width, height = shape.width, shape.height
                        
                        # Empêcher le débordement
                        if width > slide_width:
                            ratio = slide_width / width
                            width = slide_width
                            height = int(height * ratio)
                            left = 0
                        if height > slide_height:
                            ratio = slide_height / height
                            height = slide_height
                            width = int(width * ratio)
                            top = 0
                        
                        slide.shapes.add_picture(nouvelle_image_io, left, top, width, height)
                        shapes_to_delete.append(shape)
                
                # --- GESTION DES TEXTES ---
                elif shape.has_text_frame:
                    is_title = shape == slide.shapes.title
                    
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = 'Inter'
                            run.font.color.rgb = texte_blanc
                            
                            if is_title:
                                run.font.size = Pt(32)
                                run.font.bold = True
                            else:
                                if run.font.size and run.font.size > Pt(20):
                                    run.font.size = Pt(24)
                                    run.font.bold = True
                                else:
                                    run.font.size = Pt(18)
                                    run.font.bold = False

            # Nettoyage des anciennes images sur cette slide
            for shape in shapes_to_delete:
                shape._element.getparent().remove(shape._element)
            
            # 4. OPTIMISATION : Forcer le vidage de la mémoire vive après chaque slide
            gc.collect()

        output_path = "Decathlon_Optimise.pptx"
        prs.save(output_path)
        
        return FileResponse(
            output_path, 
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", 
            filename="Decathlon_Optimise.pptx"
        )
        
    except Exception as e:
        print(f"Erreur globale : {e}")
        # On renvoie une erreur propre au frontend
        raise HTTPException(status_code=500, detail="Le fichier est trop lourd ou corrompu pour le serveur.")
        
    finally:
        # On s'assure de toujours supprimer le fichier d'entrée, même en cas de crash
        if os.path.exists(input_path):
            os.remove(input_path)

@app.get("/")
async def afficher_page_accueil():
    with open("index.html", "r", encoding="utf-8") as f:
        return HTMLResponse(content=f.read())
